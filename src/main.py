"""
Generate a PowerPoint slide from hierarchy Excel data.

Drawing strategy: all connectors are drawn first, then all circles, so that
connectors always appear behind circles regardless of layout mode.
"""

import math
import yaml
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


ROOT_DIR = Path(__file__).resolve().parent.parent


def load_config():
    config_path = ROOT_DIR / "config" / "config.yaml"
    with open(config_path, "r") as f:
        return yaml.safe_load(f)


def hex_to_rgb(hex_str):
    return RGBColor(
        int(hex_str[0:2], 16),
        int(hex_str[2:4], 16),
        int(hex_str[4:6], 16),
    )


def read_sheet_mapping(wb, sheet_name, header_row, parent_col, child_col):
    """Read a sheet and return an ordered dict of parent -> [children]."""
    ws = wb[sheet_name]
    mapping = {}
    for row in range(header_row + 1, ws.max_row + 1):
        parent = ws[f"{parent_col}{row}"].value
        child = ws[f"{child_col}{row}"].value
        if parent and child:
            mapping.setdefault(parent, []).append(child)
    return mapping


def read_components(wb, config):
    """Read the Components sheet and return a dict of (type, name) -> progress.

    Type values are normalised: 'L1', 'Agent', 'Tool' (stripped/title-cased).
    """
    comp_cfg = config["sheets"]["Components"]
    ws = wb["Components"]
    header_row = comp_cfg["header_row"]
    type_col = comp_cfg["type_column"]
    name_col = comp_cfg["name_column"]
    prog_col = comp_cfg["progress_column"]

    progress = {}
    for row in range(header_row + 1, ws.max_row + 1):
        ctype = ws[f"{type_col}{row}"].value
        cname = ws[f"{name_col}{row}"].value
        cprog = ws[f"{prog_col}{row}"].value
        if ctype and cname and cprog is not None:
            progress[(ctype.strip(), cname.strip())] = int(cprog)
    return progress


def read_all_data(config):
    """Read L12A, A2T, and Components sheets."""
    wb = openpyxl.load_workbook(ROOT_DIR / config["input_file"], data_only=True)

    l12a_cfg = config["sheets"]["L12A"]
    l1_to_agents = read_sheet_mapping(
        wb, "L12A", l12a_cfg["header_row"],
        l12a_cfg["l1_column"], l12a_cfg["agents_column"],
    )

    a2t_cfg = config["sheets"]["A2T"]
    agent_to_tools = read_sheet_mapping(
        wb, "A2T", a2t_cfg["header_row"],
        a2t_cfg["agents_column"], a2t_cfg["tools_column"],
    )

    component_progress = read_components(wb, config)

    wb.close()

    # Only keep agents in A2T that also appear in L12A
    known_agents = {a for agents in l1_to_agents.values() for a in agents}
    agent_to_tools = {
        agent: tools for agent, tools in agent_to_tools.items()
        if agent in known_agents
    }

    return l1_to_agents, agent_to_tools, component_progress


# ---------------------------------------------------------------------------
# Low-level drawing primitives
# ---------------------------------------------------------------------------

def draw_line(slide, x1, y1, x2, y2, color, weight):
    """Add a straight connector line between two points."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = weight


def draw_circle(slide, center_x, center_y, size_emu, fill_color, font_size,
                font_color, text, abbrev_chars=0, label_font_size=None,
                label_font_color=None):
    """Add a circle centred at (center_x, center_y) with optional label below."""
    half = size_emu // 2
    shape = slide.shapes.add_shape(
        9, center_x - half, center_y - half, size_emu, size_emu,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    needs_label = abbrev_chars > 0 and len(text) > abbrev_chars
    inside_text = text[:abbrev_chars] if needs_label else text

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = inside_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = font_size
    p.font.color.rgb = font_color

    if needs_label:
        label_width = Inches(1.5)
        label_height = Inches(0.3)
        label_left = center_x - label_width // 2
        label_top = center_y + half + Inches(0.02)
        tb = slide.shapes.add_textbox(label_left, label_top, label_width, label_height)
        tf_label = tb.text_frame
        tf_label.word_wrap = True
        p_label = tf_label.paragraphs[0]
        p_label.text = text
        p_label.alignment = PP_ALIGN.CENTER
        p_label.font.size = label_font_size or font_size
        p_label.font.color.rgb = label_font_color or font_color


HARVEY_BALL_CHARS = {
    0: "\u25CB",  # ○ empty
    1: "\u25D4",  # ◔ quarter
    2: "\u25D1",  # ◑ half
    3: "\u25D5",  # ◕ three-quarter
    4: "\u25CF",  # ● full
}


def draw_harvey_ball(slide, center_x, center_y, size_emu, progress,
                     outline_color, outline_weight, fill_color, bg_color):
    """Draw a Harvey ball at (center_x, center_y) using a Unicode character.

    progress: 0=empty, 1=quarter, 2=half, 3=three-quarter, 4=full.
    """
    # Use a text box with the Unicode Harvey ball character
    # Make the box slightly larger than the desired size for text fitting
    box_size = int(size_emu * 2.0)
    left = center_x - box_size // 2
    top = center_y - box_size // 2

    tb = slide.shapes.add_textbox(left, top, box_size, box_size)
    tf = tb.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.text = HARVEY_BALL_CHARS.get(progress, HARVEY_BALL_CHARS[0])
    p.alignment = PP_ALIGN.CENTER
    # Scale font to match desired harvey ball size
    p.font.size = int(size_emu * 0.85)
    p.font.color.rgb = fill_color


def _make_circle_style(cfg):
    """Build a style dict from a circle config section."""
    style = {
        "size": Inches(cfg["size_inches"]),
        "fill_color": hex_to_rgb(cfg["fill_color"]),
        "font_size": Pt(cfg["font_size"]),
        "font_color": hex_to_rgb(cfg.get("font_color", "000000")),
        "abbrev_chars": cfg.get("abbrev_chars", 0),
    }
    if cfg.get("label_font_size"):
        style["label_font_size"] = Pt(cfg["label_font_size"])
    if cfg.get("label_font_color"):
        style["label_font_color"] = hex_to_rgb(cfg["label_font_color"])
    return style


# ---------------------------------------------------------------------------
# Batch drawing helpers — work on collected positions
# ---------------------------------------------------------------------------

def _draw_all_connectors(slide, connectors):
    """Draw all connector lines. Each entry: (x1, y1, x2, y2, color, weight)."""
    for x1, y1, x2, y2, color, weight in connectors:
        draw_line(slide, x1, y1, x2, y2, color, weight)


def _draw_all_circles(slide, circles):
    """Draw all circles. Each entry: (cx, cy, style, text)."""
    for cx, cy, style, text in circles:
        draw_circle(slide, cx, cy, style["size"], style["fill_color"],
                    style["font_size"], style["font_color"], text,
                    abbrev_chars=style.get("abbrev_chars", 0),
                    label_font_size=style.get("label_font_size"),
                    label_font_color=style.get("label_font_color"))


# ---------------------------------------------------------------------------
# Auto-spacing helpers
# ---------------------------------------------------------------------------

def _min_orbit_radius(n_children, child_radius, padding):
    """Minimum orbit distance so N children of given radius don't overlap.

    Children are placed equidistant on a circle.  Adjacent children are
    separated by angle 2*pi/N.  They touch when:
        2 * orbit * sin(pi/N) = 2 * child_radius + padding
    So:
        orbit = (child_radius + padding/2) / sin(pi/N)
    For N=1, no sibling overlap is possible, return child_radius + padding.
    """
    if n_children <= 1:
        return child_radius + padding
    return (child_radius + padding / 2) / math.sin(math.pi / n_children)


def _safe_orbit_distance(n_children, child_radius, configured_dist, padding):
    """Return the larger of configured distance and the minimum safe distance."""
    min_dist = _min_orbit_radius(n_children, child_radius, padding)
    return max(configured_dist, min_dist)


# ---------------------------------------------------------------------------
# Position calculation — no drawing, just returns positions + draw lists
# ---------------------------------------------------------------------------

def _calc_surrounding_positions(parent_to_children, parent_centers, center_dist,
                                child_radius=0, padding=0):
    """Calculate child positions orbiting around parent centers.
    If child_radius > 0, auto-expands orbit to prevent sibling overlap.
    Returns dict of child_name -> (cx, cy).
    """
    child_centers = {}
    for parent_name, children in parent_to_children.items():
        if parent_name not in parent_centers:
            continue
        px, py = parent_centers[parent_name]
        n = len(children)
        dist = _safe_orbit_distance(n, child_radius, center_dist, padding) if child_radius else center_dist
        for j, child_name in enumerate(children):
            angle = 2 * math.pi * j / n - math.pi / 2
            child_centers[child_name] = (
                px + int(dist * math.cos(angle)),
                py + int(dist * math.sin(angle)),
            )
    return child_centers


def _calc_row_positions(parent_to_children, row_y, child_size, slide_width):
    """Calculate child positions in a centered horizontal row.
    Returns dict of child_name -> (cx, cy) and list of (parent_name, child_name).
    """
    all_children = []
    for parent_name, children in parent_to_children.items():
        for child_name in children:
            all_children.append((parent_name, child_name))

    if not all_children:
        return {}, []

    padding = Inches(0.3)
    total_width = len(all_children) * child_size + (len(all_children) - 1) * padding
    start_x = (slide_width - total_width) // 2 + child_size // 2

    child_centers = {}
    for i, (parent_name, child_name) in enumerate(all_children):
        cx = start_x + i * (child_size + padding)
        child_centers[child_name] = (cx, row_y)

    return child_centers, all_children


def _calc_l1_agents_surrounding(l1_to_agents, sheet_cfg, agent_size, padding=0,
                                top_margin=None):
    """Calculate L1 and agent positions for surrounding layout.
    Auto-expands orbit distances per cluster and spaces clusters apart.
    Returns l1_centers, agent_centers.
    """
    configured_dist = Inches(sheet_cfg["center_distance_inches"])
    agent_radius = agent_size // 2
    cluster_padding = Inches(0.3)
    start_x = Inches(1.0)
    center_y = top_margin + Inches(2.5) if top_margin is not None else Inches(3.0)

    l1_centers = {}
    agent_centers = {}

    # First pass: compute per-cluster orbit radius
    cluster_radii = []
    for l1_name, agents in l1_to_agents.items():
        n = len(agents)
        dist = _safe_orbit_distance(n, agent_radius, configured_dist, padding)
        cluster_radii.append(dist + agent_radius)

    # Place clusters so they don't overlap
    cx = start_x + cluster_radii[0]
    for i, (l1_name, agents) in enumerate(l1_to_agents.items()):
        if i > 0:
            cx += cluster_radii[i] + cluster_padding

        cy = center_y
        l1_centers[l1_name] = (int(cx), cy)

        n = len(agents)
        dist = _safe_orbit_distance(n, agent_radius, configured_dist, padding)
        for j, agent_name in enumerate(agents):
            angle = 2 * math.pi * j / n - math.pi / 2
            agent_centers[agent_name] = (
                int(cx + dist * math.cos(angle)),
                int(cy + dist * math.sin(angle)),
            )

        if i < len(cluster_radii) - 1:
            cx += cluster_radii[i] + cluster_padding

    return l1_centers, agent_centers


def _calc_l1_agents_bottom(l1_to_agents, l1_size, agent_size, top_margin=None):
    """Calculate L1 and agent positions for bottom layout.
    Returns l1_centers, agent_centers, agent_parent_list.
    """
    slide_width = Inches(10.0)
    margin = top_margin if top_margin is not None else Inches(0.5)

    # L1 row
    l1_names = list(l1_to_agents.keys())
    l1_padding = Inches(0.5)
    total_l1_width = len(l1_names) * l1_size + (len(l1_names) - 1) * l1_padding
    l1_start_x = (slide_width - total_l1_width) // 2 + l1_size // 2
    l1_y = margin + Inches(1.0)

    l1_centers = {}
    for i, name in enumerate(l1_names):
        l1_centers[name] = (l1_start_x + i * (l1_size + l1_padding), l1_y)

    # Agent row
    all_agents = []
    for l1_name, agents in l1_to_agents.items():
        for agent_name in agents:
            all_agents.append((l1_name, agent_name))

    agent_padding = Inches(0.3)
    total_agent_width = len(all_agents) * agent_size + (len(all_agents) - 1) * agent_padding
    agent_start_x = (slide_width - total_agent_width) // 2 + agent_size // 2
    agent_y = margin + Inches(4.0)

    agent_centers = {}
    for i, (l1_name, agent_name) in enumerate(all_agents):
        agent_centers[agent_name] = (
            agent_start_x + i * (agent_size + agent_padding), agent_y,
        )

    return l1_centers, agent_centers, all_agents


# ---------------------------------------------------------------------------
# Connector collection helpers
# ---------------------------------------------------------------------------

def _collect_connectors(parent_centers, child_centers, parent_to_children,
                        conn_color, conn_weight):
    """Build a list of connector tuples from parent->child mappings."""
    connectors = []
    for parent_name, children in parent_to_children.items():
        if parent_name not in parent_centers:
            continue
        px, py = parent_centers[parent_name]
        for child_name in children:
            if child_name in child_centers:
                cx, cy = child_centers[child_name]
                connectors.append((px, py, cx, cy, conn_color, conn_weight))
    return connectors


# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------

def _calc_harvey_position(circle_cx, circle_cy, circle_size, hb_size, offset, position):
    """Calculate Harvey ball center given the parent circle center and config."""
    r = circle_size // 2
    hb_r = hb_size // 2
    dist = r + offset + hb_r  # from circle center to harvey ball center

    positions = {
        "bottom-right": (math.pi / 4),       # 45 degrees
        "bottom-left":  (3 * math.pi / 4),   # 135 degrees
        "top-left":     (5 * math.pi / 4),   # 225 degrees
        "top-right":    (7 * math.pi / 4),   # 315 degrees
        "right":        (0,),
        "bottom":       (math.pi / 2,),
        "left":         (math.pi,),
        "top":          (3 * math.pi / 2,),
    }

    angle = positions.get(position, math.pi / 4)  # default bottom-right
    if isinstance(angle, tuple):
        angle = angle[0]

    hx = circle_cx + int(dist * math.cos(angle))
    hy = circle_cy + int(dist * math.sin(angle))
    return hx, hy


def _draw_all_harvey_balls(slide, harvey_balls):
    """Draw all Harvey balls. Each entry: (cx, cy, size, progress, outline_color,
    outline_weight, fill_color, bg_color)."""
    for cx, cy, size, progress, oc, ow, fc, bg in harvey_balls:
        draw_harvey_ball(slide, cx, cy, size, progress, oc, ow, fc, bg)


def build_slide(config, l1_to_agents, agent_to_tools, component_progress):
    """Create a PowerPoint with circles for L1, agents, and tools.

    All connectors are drawn first, then all circles, ensuring connectors
    are always behind circles.
    """
    l12a_cfg = config["sheets"]["L12A"]
    l12a_conn = l12a_cfg["connector"]
    l1_style = _make_circle_style(l12a_cfg["l1_circle"])
    agent_style = _make_circle_style(l12a_cfg["agent_circle"])
    l12a_draw_conn = l12a_conn.get("enabled", True)
    l12a_conn_color = hex_to_rgb(l12a_conn["color"])
    l12a_conn_weight = Pt(l12a_conn["weight_pt"])

    a2t_cfg = config["sheets"]["A2T"]
    a2t_conn = a2t_cfg["connector"]
    tool_style = _make_circle_style(a2t_cfg["tool_circle"])
    a2t_draw_conn = a2t_conn.get("enabled", True)
    a2t_conn_color = hex_to_rgb(a2t_conn["color"])
    a2t_conn_weight = Pt(a2t_conn["weight_pt"])

    slide_width = Inches(10.0)

    # Auto-spacing config
    auto_cfg = config.get("auto_spacing", {})
    auto_enabled = auto_cfg.get("enabled", False)
    auto_padding = Inches(auto_cfg.get("padding_inches", 0.05)) if auto_enabled else 0

    # Slide layout config
    slide_cfg = config.get("slide", {})
    top_margin = Inches(slide_cfg.get("top_margin_inches", 0.5))

    # ---- Phase 1: Calculate all positions ----

    agent_layout = l12a_cfg.get("agent_layout", "bottom")
    if agent_layout == "surrounding":
        l1_centers, agent_centers = _calc_l1_agents_surrounding(
            l1_to_agents, l12a_cfg, agent_style["size"],
            padding=auto_padding, top_margin=top_margin)
    else:
        l1_centers, agent_centers, _ = _calc_l1_agents_bottom(
            l1_to_agents, l1_style["size"], agent_style["size"],
            top_margin=top_margin)

    tool_layout = a2t_cfg.get("tool_layout", "top")
    if tool_layout == "surrounding":
        center_dist = Inches(a2t_cfg["center_distance_inches"])
        tool_centers = _calc_surrounding_positions(
            agent_to_tools, agent_centers, center_dist,
            child_radius=tool_style["size"] // 2, padding=auto_padding)
    else:
        tool_centers, _ = _calc_row_positions(
            agent_to_tools, row_y=top_margin,
            child_size=tool_style["size"], slide_width=slide_width)

    # ---- Phase 2: Collect all connectors ----

    all_connectors = []
    if l12a_draw_conn:
        all_connectors += _collect_connectors(
            l1_centers, agent_centers, l1_to_agents,
            l12a_conn_color, l12a_conn_weight)
    if a2t_draw_conn:
        all_connectors += _collect_connectors(
            agent_centers, tool_centers, agent_to_tools,
            a2t_conn_color, a2t_conn_weight)

    # ---- Phase 3: Collect all circles ----

    all_circles = []
    for name, (cx, cy) in l1_centers.items():
        all_circles.append((cx, cy, l1_style, name))
    for name, (cx, cy) in agent_centers.items():
        all_circles.append((cx, cy, agent_style, name))
    for name, (cx, cy) in tool_centers.items():
        all_circles.append((cx, cy, tool_style, name))

    # ---- Phase 4: Collect Harvey balls ----

    all_harvey_balls = []
    comp_cfg = config["sheets"].get("Components", {})
    hb_cfg = comp_cfg.get("harvey_ball", {})

    if hb_cfg and hb_cfg.get("enabled", True) and component_progress:
        hb_size = Inches(hb_cfg.get("size_inches", 0.15))
        hb_offset = Inches(hb_cfg.get("offset_inches", 0.08))
        hb_outline = hex_to_rgb(hb_cfg.get("outline_color", "000000"))
        hb_weight = Pt(hb_cfg.get("outline_weight_pt", 0.75))
        hb_fill = hex_to_rgb(hb_cfg.get("fill_color", "2E75B6"))
        hb_bg = hex_to_rgb(hb_cfg.get("background_color", "FFFFFF"))
        hb_position = hb_cfg.get("position", "bottom-right")

        # Map component types to their center dicts and circle sizes
        type_map = {
            "L1":    (l1_centers, l1_style["size"]),
            "Agent": (agent_centers, agent_style["size"]),
            "Tool":  (tool_centers, tool_style["size"]),
        }

        for (ctype, cname), progress in component_progress.items():
            if ctype not in type_map:
                continue
            centers, circle_size = type_map[ctype]
            if cname not in centers:
                continue
            cx, cy = centers[cname]
            hx, hy = _calc_harvey_position(cx, cy, circle_size, hb_size,
                                           hb_offset, hb_position)
            all_harvey_balls.append(
                (hx, hy, hb_size, progress, hb_outline, hb_weight, hb_fill, hb_bg))

    # ---- Phase 5: Draw — connectors first, then circles, then Harvey balls ----

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _draw_all_connectors(slide, all_connectors)
    _draw_all_circles(slide, all_circles)
    _draw_all_harvey_balls(slide, all_harvey_balls)

    output_path = ROOT_DIR / config["output_file"]
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    config = load_config()
    l1_to_agents, agent_to_tools, component_progress = read_all_data(config)

    print("L1 -> Agents:")
    for l1, agents in l1_to_agents.items():
        print(f"  {l1}: {agents}")

    print("Agent -> Tools:")
    for agent, tools in agent_to_tools.items():
        print(f"  {agent}: {tools}")

    print("Component Progress:")
    for (ctype, cname), prog in component_progress.items():
        print(f"  {ctype}/{cname}: {prog}")

    output = build_slide(config, l1_to_agents, agent_to_tools, component_progress)
    print(f"Saved to {output}")


if __name__ == "__main__":
    main()
